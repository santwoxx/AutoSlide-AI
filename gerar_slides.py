#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script de Automação gerado por Automatizador de Slides.
Este script monta uma apresentação de slides profissional com layouts modernos usando a biblioteca 'python-pptx'.

Requisitos:
    pip install python-pptx

Modo de Uso:
    python gerar_slides.py
"""

import os
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- PALETA DE CORES (Tema: Corporativo Royal (Azul, Branco & Ouro)) ---
COLOR_BG = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(15, 23, 42)
COLOR_ACCENT = RGBColor(212, 175, 55)
COLOR_CARD = RGBColor(239, 246, 255)

FONT_TITLE = 'Arial'
FONT_BODY = 'Arial'

def create_deck():
    # Inicializa apresentação widescreen 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Vamos usar o layout embranquecido/vazio padrão para controle total dos elementos
    blank_layout = prs.slide_layouts[6]

    # --- DICIONÁRIO DE IMAGENS EMBUTIDAS (BASE64) ---
    # As imagens inseridas no editor web são salvas aqui em formato compacto 
    # para auto-suficiência completa do script python localmente!
    images_base64 = {
    }

    # Salva as imagens temporárias decodificadas
    temp_files = []
    print("Decodificando imagens locais...")
    for img_id, b64_data in images_base64.items():
        try:
            filename = f"temp_image_{img_id}.png"
            with open(filename, "wb") as f:
                f.write(base64.b64decode(b64_data.strip()))
            temp_files.append(filename)
        except Exception as e:
            print(f"Erro ao salvar imagem {img_id}: {e}")

    # --- AJUDANTES DE MONTAGEM ---
    
    def apply_solid_background(slide, color):
        """Define a cor de fundo sólida do slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_card_shape(slide, left, top, width, height, fill_color, line_color=None):
        """Cria um retângulo decorativo de fundo / fundo de conteúdo."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background() # Sem borda
        return shape

    def format_text_frame(tf):
        """Elimina as margens internas padrão de caixas de textos."""
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

    print("Montando slides da apresentação...")

    # --- SLIDE 1 (title) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide, COLOR_BG)
    # Layout: Capa da Apresentação
    # Caixa de texto centralizada para o título principal e subtítulo
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
    tf = tx_box.text_frame
    format_text_frame(tf)
    
    # Título Principal
    p_title = tf.paragraphs[0]
    p_title.text = "AUTOMAÇÃO DE SLIDES PRO"
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT
    p_title.alignment = PP_ALIGN.CENTER
    p_title.space_after = Pt(24)
    
    # Subtítulo (se houver)
    p_sub = tf.add_paragraph()
    p_sub.text = "Gere apresentações profissionais e scripts python-pptx dinâmicos de forma mecânica e sem IA"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_after = Pt(18)
    # Rodapé / Apresentador
    p_footer = tf.add_paragraph()
    p_footer.text = "Gerado via Automação Python-PPTX  •  16:9 Widescreen"
    p_footer.font.name = FONT_BODY
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = COLOR_TEXT
    p_footer.alignment = PP_ALIGN.CENTER

    # --- SLIDE 2 (bullets) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide, COLOR_BG)
    # Layout: Lista de Pontos Importantes (Bullets)
    # Título do slide
    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
    tf_t = tx_title.text_frame
    format_text_frame(tf_t)
    p_t = tf_t.paragraphs[0]
    p_t.text = "Por que usar Automação Estruturada?"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_ACCENT
    
    # Caixa com Bullet Points
    tx_bullets = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
    tf_b = tx_bullets.text_frame
    format_text_frame(tf_b)
    p_b0 = tf_b.paragraphs[0]
    p_b0.text = "• Garantia de consistência visual completa entre múltiplos decks de slides"
    p_b0.font.name = FONT_BODY
    p_b0.font.size = Pt(18)
    p_b0.font.color.rgb = COLOR_TEXT
    p_b0.space_after = Pt(14)
    p_bx = tf_b.add_paragraph()
    p_bx.text = "• Velocidade absurda: monte apresentações de 20+ slides em segundos"
    p_bx.font.name = FONT_BODY
    p_bx.font.size = Pt(18)
    p_bx.font.color.rgb = COLOR_TEXT
    p_bx.space_after = Pt(14)
    p_bx = tf_b.add_paragraph()
    p_bx.text = "• Independência de modelos de IA e indisponibilidades de redes externas"
    p_bx.font.name = FONT_BODY
    p_bx.font.size = Pt(18)
    p_bx.font.color.rgb = COLOR_TEXT
    p_bx.space_after = Pt(14)
    p_bx = tf_b.add_paragraph()
    p_bx.text = "• Integração facilitada com cronogramas, bancos de dados locais ou logs de sistemas"
    p_bx.font.name = FONT_BODY
    p_bx.font.size = Pt(18)
    p_bx.font.color.rgb = COLOR_TEXT
    p_bx.space_after = Pt(14)
    p_bx = tf_b.add_paragraph()
    p_bx.text = "• Código Python gerado totalmente portável e customizável localmente"
    p_bx.font.name = FONT_BODY
    p_bx.font.size = Pt(18)
    p_bx.font.color.rgb = COLOR_TEXT
    p_bx.space_after = Pt(14)

    # --- SLIDE 3 (comparison) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide, COLOR_BG)
    # Layout: Comparativo em 2 colunas
    # Título do slide
    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
    tf_t = tx_title.text_frame
    format_text_frame(tf_t)
    p_t = tf_t.paragraphs[0]
    p_t.text = "Comparativo: Web UI vs Script Local"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_ACCENT

    # Coluna 1 (Esquerda) - Card visual
    add_card_shape(slide, Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3), COLOR_CARD)
    tx_col1 = slide.shapes.add_textbox(Inches(1.3), Inches(2.5), Inches(4.7), Inches(3.7))
    tf_col1 = tx_col1.text_frame
    format_text_frame(tf_col1)
    
    p_c1_t = tf_col1.paragraphs[0]
    p_c1_t.text = "Visualização Web Dinâmica"
    p_c1_t.font.name = FONT_TITLE
    p_c1_t.font.size = Pt(22)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = COLOR_ACCENT
    p_c1_t.space_after = Pt(10)
    
    p_c1_txt = tf_col1.add_paragraph()
    p_c1_txt.text = "Edite em tempo real, gerencie imagens por drag & drop, assista animações fluidas e exporte apresentações formatadas (.pptx) em um único clique no navegador."
    p_c1_txt.font.name = FONT_BODY
    p_c1_txt.font.size = Pt(14)
    p_c1_txt.font.color.rgb = COLOR_TEXT

    # Coluna 2 (Direita) - Card visual
    add_card_shape(slide, Inches(6.8), Inches(2.2), Inches(5.3), Inches(4.3), COLOR_CARD)
    tx_col2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(4.7), Inches(3.7))
    tf_col2 = tx_col2.text_frame
    format_text_frame(tf_col2)
    
    p_c2_t = tf_col2.paragraphs[0]
    p_c2_t.text = "Automação Local em Lote"
    p_c2_t.font.name = FONT_TITLE
    p_c2_t.font.size = Pt(22)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = COLOR_ACCENT
    p_c2_t.space_after = Pt(10)
    
    p_c2_txt = tf_col2.add_paragraph()
    p_c2_txt.text = "Rode o script Python autônomo. Ideal para inserir em fluxos integrados e automatizar relatórios de servidores, PDF wrappers ou cronjobs industriais."
    p_c2_txt.font.name = FONT_BODY
    p_c2_txt.font.size = Pt(14)
    p_c2_txt.font.color.rgb = COLOR_TEXT

    # --- SLIDE 4 (metrics) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide, COLOR_BG)
    # Layout: Card de Métricas / Estatísticas
    # Título do slide
    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
    tf_t = tx_title.text_frame
    format_text_frame(tf_t)
    p_t = tf_t.paragraphs[0]
    p_t.text = "Impacto de Produtividade"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_ACCENT

    # Caixa Gigante para o Número
    tx_num = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.2))
    tf_num = tx_num.text_frame
    format_text_frame(tf_num)
    p_num = tf_num.paragraphs[0]
    p_num.text = "10x Mais Rápido"
    p_num.font.name = FONT_TITLE
    p_num.font.bold = True
    p_num.font.size = Pt(96)
    p_num.font.color.rgb = COLOR_ACCENT
    p_num.alignment = PP_ALIGN.CENTER

    # Detalhar Métrica / Label
    tx_det = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.333), Inches(1.8))
    tf_det = tx_det.text_frame
    format_text_frame(tf_det)
    p_lbl = tf_det.paragraphs[0]
    p_lbl.text = "Economia com Geração em Lote"
    p_lbl.font.name = FONT_BODY
    p_lbl.font.bold = True
    p_lbl.font.size = Pt(24)
    p_lbl.font.color.rgb = COLOR_TEXT
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.space_after = Pt(8)
    
    p_subt = tf_det.add_paragraph()
    p_subt.text = "Substituição de processos de montagem manual repetitiva por scripts deterministas estáveis e sem erro humano."
    p_subt.font.name = FONT_BODY
    p_subt.font.size = Pt(16)
    p_subt.font.color.rgb = COLOR_TEXT
    p_subt.alignment = PP_ALIGN.CENTER

    # --- SLIDE 5 (quote) ---
    slide = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide, COLOR_BG)
    # Layout: Depoimento / Citação Editorial (Estilo Minimalista)
    # Card com bordas sutis ou fundo diferente
    add_card_shape(slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.5), COLOR_CARD)
    
    tx_qt = slide.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(3.7))
    tf_qt = tx_qt.text_frame
    format_text_frame(tf_qt)
    
    p_qt = tf_qt.paragraphs[0]
    p_qt.text = '"A verdadeira eficiência não está em reescrever layouts do zero todos os dias, mas sim em parametrizá-los para focar exclusivamente na qualidade do conteúdo."'
    p_qt.font.name = FONT_TITLE
    p_qt.font.size = Pt(24)
    p_qt.font.italic = True
    p_qt.font.color.rgb = COLOR_TEXT
    p_qt.alignment = PP_ALIGN.CENTER
    p_qt.space_after = Pt(24)
    
    p_auth = tf_qt.add_paragraph()
    p_auth.text = "— Linus Torvalds adaptado"
    p_auth.font.name = FONT_BODY
    p_auth.font.bold = True
    p_auth.font.size = Pt(18)
    p_auth.font.color.rgb = COLOR_ACCENT
    p_auth.alignment = PP_ALIGN.CENTER
    p_auth.space_after = Pt(4)
    p_autht = tf_qt.add_paragraph()
    p_autht.text = "Entusiasta de Automações Eficientes"
    p_autht.font.name = FONT_BODY
    p_autht.font.size = Pt(13)
    p_autht.font.color.rgb = COLOR_TEXT
    p_autht.alignment = PP_ALIGN.CENTER

    # --- SALVAR O ARQUIVO FINAL ---
    out_filename = "apresentacao_automatizada.pptx"
    prs.save(out_filename)
    print(f"\nSucesso! Apresentação salva como '{out_filename}'.")

    # Limpeza de arquivos temporários
    print("Limpando arquivos temporários decodificados...")
    for filename in temp_files:
        try:
            if os.path.exists(filename):
                os.remove(filename)
        except Exception as e:
            pass
            
    print("Processo finalizado!")

if __name__ == "__main__":
    create_deck()
